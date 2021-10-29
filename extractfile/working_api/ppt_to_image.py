#todo require some work
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from coutoEditor.global_variable import BASE_URL,BASE_DIR
from .extract_desc import get_desc

def ppt_to_png(file_path,file_name,ext):
    #covert pptslides into image
    path = file_path.replace(ext,'temp')

    dir = BASE_DIR + "/media/extract_file/ppt_file/"
    pdf_path = dir + file_name + ".pdf"
    jpg_path = dir + file_name + ".jpg"

    os.system("soffice --headless --convert-to pdf " + file_path +
              " --outdir " +
              dir)

    os.system('convert -density 150 '+
              pdf_path +
              ' -quality 80 '+
              jpg_path)

    return path

def get_pptx_data(file_path):
    prs=Presentation(file_path)
    l=[]
    for slide in prs.slides:
        s=''
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                s+=shape.text
                s = s.replace('\n' , '')
                s = ' '.join(s.split())
        l.append(s)
    return l

def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape

def extract_image_from_ppt(filepath):
    for picture in iter_picture_shapes(Presentation(filepath)):
        image = picture.image
        # ---get image "file" contents---
        image_bytes = image.blob
        # ---make up a name for the file, e.g. 'image.jpg'---
        image_filename = 'image.b%s' % image.ext
        with open(image_filename, 'wb') as f:
            f.write(image_bytes)

def extract_info_from_ppt(file_path,file_name):
    #extract text from every slide
    ppt_text = get_pptx_data(file_path)

    ext = "ppt" if file_path.endswith('ppt') else "pptx"
    file_name = file_name.replace("." + ext, "")

    #title and description from slide
    desc_text = ''
    for i in range(len(ppt_text)):
        desc_text  += ppt_text[i] + "\n"
    desc = get_desc(desc_text)
    ppt_data = {}
    ppt_data['title'] = None
    ppt_data['desc'] = desc

    #convert slides into ppt
    path = ppt_to_png(file_path,file_name,ext)

    #extract image form ppt
    # extract_image_from_ppt(filepath)

    ppt_data['images'] = []
    slides = {}
    dir = "media/extract_file/ppt_file/"
    path = os.path.join(BASE_URL,dir,file_name)
    for i in range(len(ppt_text)):
        slides[i] = {"img_path": path +"-"+ str(i) + '.jpg', "text": ppt_text[i]}
    ppt_data["slides"] = slides

    return ppt_data
